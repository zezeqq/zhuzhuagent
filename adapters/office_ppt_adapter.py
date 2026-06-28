from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from core.app_identity import APP_NAME

from artifacts.artifact_manager import ensure_export_dir

_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
_ACCENT_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
_MEDIUM_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
_SLIDE_BG = RGBColor(0x0F, 0x17, 0x2A)
_CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
_GRADIENT_START = RGBColor(0x1E, 0x3A, 0x5F)
_GRADIENT_END = RGBColor(0x0F, 0x17, 0x2A)

SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)


def _set_slide_bg(slide, color: RGBColor = _SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_rect(slide, left, top, width, height, fill_color: RGBColor, corner_radius=Emu(100000)):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    sp = shape._element
    prstGeom = sp.find(qn("a:prstGeom"), sp.nsmap) if hasattr(sp, 'nsmap') else None
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=Pt(18),
                  color=_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _build_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _SLIDE_BG)

    accent_bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _ACCENT
    accent_bar.line.fill.background()

    _add_text_box(slide, Inches(1.5), Inches(2.0), Inches(7), Inches(1.5),
                  title, font_size=Pt(40), color=_WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    sub_text = subtitle or f"{APP_NAME} 自动生成"
    _add_text_box(slide, Inches(1.5), Inches(3.6), Inches(7), Inches(0.8),
                  sub_text, font_size=Pt(18), color=_MEDIUM_GRAY,
                  alignment=PP_ALIGN.CENTER)

    left_bar = slide.shapes.add_shape(1, Inches(4.2), Inches(3.3), Inches(1.6), Inches(0.04))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = _ACCENT
    left_bar.line.fill.background()


def _build_content_slide(prs, slide_title: str, bullets: list[str], slide_num: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _SLIDE_BG)

    header_bg = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = _DARK_BG
    header_bg.line.fill.background()

    accent_line = slide.shapes.add_shape(1, Emu(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = _ACCENT
    accent_line.line.fill.background()

    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.7),
                  slide_title, font_size=Pt(28), color=_WHITE, bold=True)

    content_top = Inches(1.6)
    content_left = Inches(0.8)
    content_width = Inches(8.4)
    line_height = Inches(0.55)

    for i, bullet in enumerate(bullets):
        y = content_top + line_height * i

        is_sub = bullet.startswith("  ") or bullet.startswith("\t")
        clean = bullet.strip().lstrip("-•·").strip()

        if is_sub:
            indent = Inches(0.5)
            marker = "›"
            fs = Pt(16)
            clr = _LIGHT_GRAY
        else:
            indent = Inches(0)
            marker = "•"
            fs = Pt(18)
            clr = _WHITE

        dot_box = _add_text_box(slide, content_left + indent, y,
                                Inches(0.3), Inches(0.45),
                                marker, font_size=Pt(18), color=_ACCENT)

        _add_text_box(slide, content_left + indent + Inches(0.35), y,
                      content_width - indent - Inches(0.35), Inches(0.45),
                      clean, font_size=fs, color=clr)

    _add_text_box(slide, Inches(8.8), Inches(7.0), Inches(1.0), Inches(0.3),
                  f"{slide_num}/{total}", font_size=Pt(10), color=_MEDIUM_GRAY,
                  alignment=PP_ALIGN.RIGHT)


def _build_end_slide(prs, title: str = "谢谢观看"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _SLIDE_BG)

    _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(7), Inches(1.2),
                  title, font_size=Pt(36), color=_WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(4.0), Inches(7), Inches(0.6),
                  f"Powered by {APP_NAME}", font_size=Pt(14),
                  color=_MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    bar = slide.shapes.add_shape(1, Inches(4.2), Inches(3.8), Inches(1.6), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()


def create_presentation(title: str, slides: list[tuple[str, list[str]]], output_name: str) -> Path:
    out_dir = ensure_export_dir()
    path = out_dir / output_name

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_title_slide(prs, title)

    total = len(slides)
    for idx, (slide_title, bullets) in enumerate(slides, 1):
        _build_content_slide(prs, slide_title, bullets, idx, total)

    _build_end_slide(prs)

    prs.save(path)
    return path
