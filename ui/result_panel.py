from __future__ import annotations

import json
import os
from pathlib import Path

from PySide6.QtCore import Qt, Signal
from PySide6.QtWidgets import (
    QAbstractItemView, QCheckBox, QFrame, QHBoxLayout, QLabel, QMessageBox,
    QPlainTextEdit, QPushButton, QScrollArea, QSizePolicy, QStackedWidget, QTreeWidget,
    QTreeWidgetItem, QVBoxLayout, QWidget,
)

from artifacts.artifact_manager import list_artifacts, remove_artifact_by_path
from db.database import query_all
from utils.file_actions import exec_file_context_menu, open_file_path, reveal_in_explorer
from utils.path_utils import exports_dir
from ui.widgets.batch_action_bar import BatchActionBar

_FILE_PATH_ROLE = Qt.UserRole
_FILE_LOADED_ROLE = Qt.UserRole + 1
_FILE_IGNORED = frozenset({".git", ".idea", ".venv", "venv", "__pycache__", "node_modules"})


class _FileTabBar(QFrame):
    """Horizontal bar of recently-clicked file tabs."""

    file_tab_clicked = None  # set externally as a callback

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("FileTabBar")
        self.setFixedHeight(32)
        self._layout = QHBoxLayout(self)
        self._layout.setContentsMargins(4, 2, 4, 2)
        self._layout.setSpacing(2)
        self._layout.addStretch()
        self._tabs: dict[str, QPushButton] = {}

    def add_tab(self, file_path: str, file_name: str) -> None:
        if file_path in self._tabs:
            self._highlight(file_path)
            return
        btn = QPushButton(file_name)
        btn.setObjectName("FileTabButton")
        btn.setCheckable(True)
        btn.setCursor(Qt.PointingHandCursor)
        btn.setFixedHeight(24)
        btn.clicked.connect(lambda _, p=file_path: self._on_click(p))
        idx = self._layout.count() - 1
        self._layout.insertWidget(idx, btn)
        self._tabs[file_path] = btn
        self._highlight(file_path)

    def _on_click(self, file_path: str) -> None:
        self._highlight(file_path)
        if self.file_tab_clicked:
            self.file_tab_clicked(file_path)

    def _highlight(self, file_path: str) -> None:
        for p, btn in self._tabs.items():
            btn.setChecked(p == file_path)


class _ArtifactCard(QFrame):
    """Artifact row: click anywhere (except buttons) to open with the OS default app."""

    open_requested = Signal(str)
    preview_requested = Signal(str)
    check_changed = Signal(str, bool)

    def __init__(self, art: dict, parent=None, on_changed=None):
        super().__init__(parent)
        self.setObjectName("ResultArtifactCard")
        self._file_path = art.get("file_path", "") or ""
        self._on_changed = on_changed
        self._multi_mode = False
        self.setCursor(Qt.PointingHandCursor)
        self.setToolTip("单击打开 · 右键更多操作")

        layout = QHBoxLayout(self)
        layout.setContentsMargins(10, 8, 10, 8)
        layout.setSpacing(8)

        self._check = QCheckBox()
        self._check.setObjectName("ArtifactCardCheck")
        self._check.setVisible(False)
        self._check.stateChanged.connect(
            lambda state: self.check_changed.emit(self._file_path, state == Qt.CheckState.Checked)
        )
        layout.addWidget(self._check)

        type_icons = {"docx": "📄", "xlsx": "📊", "pptx": "📑", "py": "🐍", "md": "📝"}
        atype = art.get("artifact_type", "")
        icon = QLabel(type_icons.get(atype, "📁"))
        icon.setFixedWidth(24)
        layout.addWidget(icon)

        info = QVBoxLayout()
        info.setSpacing(2)
        name = QLabel(art.get("artifact_name", "未命名"))
        name.setObjectName("ArtifactCardName")
        path_label = QLabel(self._file_path)
        path_label.setObjectName("MutedLabel")
        path_label.setWordWrap(True)
        info.addWidget(name)
        info.addWidget(path_label)
        layout.addLayout(info, 1)

        preview_btn = QPushButton("预览")
        preview_btn.setProperty("variant", "secondary")
        preview_btn.setCursor(Qt.PointingHandCursor)
        preview_btn.clicked.connect(lambda: self.preview_requested.emit(self._file_path))
        layout.addWidget(preview_btn)

        open_btn = QPushButton("打开")
        open_btn.setProperty("variant", "secondary")
        open_btn.setCursor(Qt.PointingHandCursor)
        open_btn.clicked.connect(lambda: self.open_requested.emit(self._file_path))
        layout.addWidget(open_btn)

    def set_multi_select_mode(self, enabled: bool) -> None:
        self._multi_mode = enabled
        self._check.setVisible(enabled)
        if not enabled:
            self._check.setChecked(False)

    def set_checked(self, checked: bool) -> None:
        self._check.blockSignals(True)
        self._check.setChecked(checked)
        self._check.blockSignals(False)

    def is_checked(self) -> bool:
        return self._check.isChecked()

    def contextMenuEvent(self, event) -> None:
        if not self._file_path:
            return
        exec_file_context_menu(
            self,
            self._file_path,
            event.globalPos(),
            on_preview=lambda: self.preview_requested.emit(self._file_path),
            on_after_delete=self._on_file_deleted,
        )
        event.accept()

    def _on_file_deleted(self) -> None:
        if self._on_changed:
            self._on_changed()

    def mouseReleaseEvent(self, event) -> None:
        if event.button() == Qt.LeftButton:
            widget = self.childAt(event.pos())
            while widget and widget is not self:
                if isinstance(widget, (QPushButton, QCheckBox)):
                    super().mouseReleaseEvent(event)
                    return
                widget = widget.parentWidget()
            if self._multi_mode:
                self._check.setChecked(not self._check.isChecked())
                return
            self.open_requested.emit(self._file_path)
        super().mouseReleaseEvent(event)


class ResultPanel(QFrame):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("ResultPanel")
        self.setMinimumWidth(340)
        self._task_id: int | None = None
        self._conversation_id: int | None = None
        self._workspace: str = ""
        self._files_stale = True
        self._artifact_cards: list[_ArtifactCard] = []
        self._artifact_multi_mode = False
        self._current_preview_path: str = ""

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        tab_bar = QFrame()
        tab_bar.setObjectName("ResultTabBar")
        tab_bar.setFixedHeight(44)
        tab_layout = QHBoxLayout(tab_bar)
        tab_layout.setContentsMargins(8, 0, 8, 0)
        tab_layout.setSpacing(0)
        self._tabs: dict[str, QPushButton] = {}
        for key, text in [("artifacts", "产物"), ("files", "文件"), ("preview", "预览"), ("changes", "变更")]:
            btn = QPushButton(text)
            btn.setObjectName("ResultTab")
            btn.setCheckable(True)
            btn.setChecked(key == "artifacts")
            btn.setCursor(Qt.PointingHandCursor)
            btn.clicked.connect(lambda _, k=key: self._switch_tab(k))
            self._tabs[key] = btn
            tab_layout.addWidget(btn)
        tab_layout.addStretch()

        collapse_btn = QPushButton("✕")
        collapse_btn.setObjectName("CollapseButton")
        collapse_btn.setFixedSize(28, 28)
        collapse_btn.setCursor(Qt.PointingHandCursor)
        collapse_btn.clicked.connect(lambda: self.setVisible(False))
        tab_layout.addWidget(collapse_btn)
        layout.addWidget(tab_bar)

        self._stack = QStackedWidget()
        self._stack.addWidget(self._build_artifacts_tab())
        self._stack.addWidget(self._build_files_tab())
        self._stack.addWidget(self._build_preview_tab())
        self._stack.addWidget(self._build_changes_tab())
        layout.addWidget(self._stack, 1)
        self._current_tab = "artifacts"

    # -- tab switching ---------------------------------------------------------

    def _switch_tab(self, key: str) -> None:
        index_map = {"artifacts": 0, "files": 1, "preview": 2, "changes": 3}
        self._stack.setCurrentIndex(index_map.get(key, 0))
        for k, btn in self._tabs.items():
            btn.setChecked(k == key)
        if self._current_tab != key:
            self._current_tab = key
            self._refresh_tab(key)
        elif key in ("artifacts", "files", "changes"):
            self._refresh_tab(key)

    def _refresh_tab(self, key: str) -> None:
        if key == "artifacts":
            self.refresh_artifacts()
        elif key == "files":
            if self._files_stale or self._file_tree.topLevelItemCount() == 0:
                self.refresh_files()
        elif key == "changes":
            self._refresh_changes()

    def refresh_current_tab(self) -> None:
        self._refresh_tab(self._current_tab)

    def refresh_file_views(self) -> None:
        """Agent 产生/修改文件后刷新产物、文件树与变更记录。"""
        self.refresh_artifacts()
        self._files_stale = True
        if self._current_tab == "files":
            self.refresh_files()
        self._refresh_changes()

    def showEvent(self, event) -> None:
        super().showEvent(event)
        self._refresh_tab(self._current_tab)

    def _build_artifacts_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)
        toolbar = QHBoxLayout()
        self._artifact_multi_btn = QPushButton("☑ 多选")
        self._artifact_multi_btn.setProperty("variant", "ghost")
        self._artifact_multi_btn.setCheckable(True)
        self._artifact_multi_btn.setCursor(Qt.PointingHandCursor)
        self._artifact_multi_btn.toggled.connect(self._toggle_artifact_multi_select)
        toolbar.addWidget(self._artifact_multi_btn)
        toolbar.addStretch()
        refresh = QPushButton("刷新")
        refresh.setProperty("variant", "ghost")
        refresh.setCursor(Qt.PointingHandCursor)
        refresh.clicked.connect(self.refresh_artifacts)
        toolbar.addWidget(refresh)
        layout.addLayout(toolbar)

        self._artifact_batch_bar = BatchActionBar()
        self._artifact_batch_bar.setVisible(False)
        self._artifact_batch_bar.select_all_clicked.connect(self._artifact_select_all)
        self._artifact_batch_bar.clear_clicked.connect(self._artifact_clear_selection)
        self._artifact_batch_bar.open_clicked.connect(self._artifact_batch_open)
        self._artifact_batch_bar.delete_clicked.connect(self._artifact_batch_delete)
        layout.addWidget(self._artifact_batch_bar)

        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        scroll.setObjectName("ArtifactScroll")
        self._artifact_container = QWidget()
        self._artifact_layout = QVBoxLayout(self._artifact_container)
        self._artifact_layout.setContentsMargins(0, 0, 0, 0)
        self._artifact_layout.setSpacing(6)
        self._artifact_layout.addStretch()
        scroll.setWidget(self._artifact_container)
        layout.addWidget(scroll, 1)
        return page

    # -- files tab -------------------------------------------------------------

    def _build_files_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)

        self._file_tab_bar = _FileTabBar()
        self._file_tab_bar.file_tab_clicked = self._on_file_tab_clicked
        layout.addWidget(self._file_tab_bar)

        toolbar = QHBoxLayout()
        self._workspace_label = QLabel("未设置工作目录")
        self._workspace_label.setObjectName("MutedLabel")
        toolbar.addWidget(self._workspace_label, 1)
        refresh = QPushButton("刷新")
        refresh.setProperty("variant", "ghost")
        refresh.setCursor(Qt.PointingHandCursor)
        refresh.clicked.connect(self.refresh_files)
        toolbar.addWidget(refresh)
        layout.addLayout(toolbar)

        self._file_batch_bar = BatchActionBar()
        self._file_batch_bar.select_all_clicked.connect(self._file_select_all)
        self._file_batch_bar.clear_clicked.connect(self._file_clear_selection)
        self._file_batch_bar.open_clicked.connect(self._file_batch_open)
        self._file_batch_bar.delete_clicked.connect(self._file_batch_delete)
        layout.addWidget(self._file_batch_bar)

        self._file_tree = QTreeWidget()
        self._file_tree.setObjectName("FileTree")
        self._file_tree.setHeaderLabels(["文件名", "大小"])
        self._file_tree.setColumnWidth(0, 200)
        self._file_tree.setSelectionMode(QAbstractItemView.ExtendedSelection)
        self._file_tree.itemSelectionChanged.connect(self._update_file_batch_bar)
        self._file_tree.itemDoubleClicked.connect(self._open_file)
        self._file_tree.itemClicked.connect(self._on_file_tree_clicked)
        self._file_tree.itemExpanded.connect(self._on_file_tree_expanded)
        self._file_tree.setContextMenuPolicy(Qt.CustomContextMenu)
        self._file_tree.customContextMenuRequested.connect(self._on_file_tree_context_menu)
        layout.addWidget(self._file_tree, 1)
        return page

    # -- preview tab -----------------------------------------------------------

    def _build_preview_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)
        self._preview_title = QLabel("选择文件以预览")
        self._preview_title.setObjectName("PreviewTitle")
        layout.addWidget(self._preview_title)
        self._preview_text = QPlainTextEdit()
        self._preview_text.setReadOnly(True)
        self._preview_text.setObjectName("PreviewText")
        layout.addWidget(self._preview_text, 1)
        toolbar = QHBoxLayout()
        toolbar.addStretch()
        open_btn = QPushButton("用默认程序打开")
        open_btn.setProperty("variant", "secondary")
        open_btn.setCursor(Qt.PointingHandCursor)
        open_btn.clicked.connect(self._open_previewed_file)
        toolbar.addWidget(open_btn)
        layout.addLayout(toolbar)
        return page

    # -- changes tab -----------------------------------------------------------

    def _build_changes_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)
        self._changes_scroll = QScrollArea()
        self._changes_scroll.setWidgetResizable(True)
        self._changes_container = QWidget()
        self._changes_layout = QVBoxLayout(self._changes_container)
        self._changes_layout.setContentsMargins(0, 0, 0, 0)
        self._changes_layout.setSpacing(6)
        self._changes_layout.addStretch()
        self._changes_scroll.setWidget(self._changes_container)
        layout.addWidget(self._changes_scroll, 1)
        return page

    # -- artifacts tab ---------------------------------------------------------

    def set_task(self, task_id: int | None) -> None:
        self._task_id = task_id
        self.refresh_file_views()

    def set_conversation(self, conv_id: int | None) -> None:
        self._conversation_id = conv_id
        self._current_preview_path = ""
        if conv_id:
            from core.task_tracker import latest_task_for_conversation
            task = latest_task_for_conversation(conv_id)
            self._task_id = task["id"] if task else None
        else:
            self._task_id = None
        self.refresh_file_views()

    # -- artifacts -------------------------------------------------------------

    def refresh_artifacts(self) -> None:
        self._artifact_cards.clear()
        while self._artifact_layout.count() > 1:
            item = self._artifact_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()
        artifacts = list_artifacts()
        if self._task_id:
            task_artifacts = [a for a in artifacts if a.get("task_id") == self._task_id]
            if task_artifacts:
                artifacts = task_artifacts
            else:
                artifacts = [a for a in artifacts if not a.get("task_id")]
        for art in artifacts:
            card = self._artifact_card(art)
            card.set_multi_select_mode(self._artifact_multi_mode)
            self._artifact_cards.append(card)
            self._artifact_layout.insertWidget(self._artifact_layout.count() - 1, card)
        self._update_artifact_batch_bar()

    def _artifact_card(self, art: dict) -> _ArtifactCard:
        card = _ArtifactCard(art, on_changed=self.refresh_file_views)
        card.open_requested.connect(self._open_file_path)
        card.preview_requested.connect(self._preview_artifact)
        card.check_changed.connect(self._on_artifact_check_changed)
        return card

    def _toggle_artifact_multi_select(self, enabled: bool) -> None:
        self._artifact_multi_mode = enabled
        self._artifact_multi_btn.setText("完成" if enabled else "☑ 多选")
        self._artifact_batch_bar.setVisible(enabled)
        for card in self._artifact_cards:
            card.set_multi_select_mode(enabled)
        if not enabled:
            self._update_artifact_batch_bar()

    def _on_artifact_check_changed(self, _path: str, _checked: bool) -> None:
        self._update_artifact_batch_bar()

    def _update_artifact_batch_bar(self) -> None:
        count = sum(1 for c in self._artifact_cards if c.is_checked())
        self._artifact_batch_bar.set_count(count)

    def _artifact_checked_paths(self) -> list[str]:
        return [c._file_path for c in self._artifact_cards if c.is_checked() and c._file_path]

    def _artifact_select_all(self) -> None:
        for card in self._artifact_cards:
            card.set_checked(True)
        self._update_artifact_batch_bar()

    def _artifact_clear_selection(self) -> None:
        for card in self._artifact_cards:
            card.set_checked(False)
        self._update_artifact_batch_bar()

    def _artifact_batch_open(self) -> None:
        for path in self._artifact_checked_paths():
            self._open_file_path(path, self)

    def _artifact_batch_delete(self) -> None:
        paths = self._artifact_checked_paths()
        if not paths:
            return
        reply = QMessageBox.question(
            self,
            "确认删除",
            f"确定删除选中的 {len(paths)} 个产物文件吗？\n此操作不可恢复。",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return
        for path in paths:
            p = Path(path)
            if p.is_file():
                try:
                    p.unlink()
                except OSError:
                    pass
            remove_artifact_by_path(path)
        self.refresh_file_views()

    @staticmethod
    def _open_file_path(file_path: str, parent: QWidget | None = None) -> bool:
        if not file_path:
            return False
        path = Path(file_path)
        if not path.exists():
            QMessageBox.warning(
                parent,
                "无法打开",
                f"文件不存在或已被移动：\n{file_path}",
            )
            return False
        try:
            os.startfile(str(path))
            return True
        except OSError as exc:
            QMessageBox.warning(
                parent,
                "无法打开",
                f"无法用系统默认程序打开该文件：\n{file_path}\n\n{exc}",
            )
            return False

    def _preview_artifact(self, file_path: str) -> None:
        """Preview an artifact file in the preview tab with format-specific handling."""
        if not file_path:
            return
        path = Path(file_path)
        if not path.exists():
            self._preview_title.setText(path.name)
            self._preview_text.setPlainText("文件不存在。")
            self._switch_tab("preview")
            return

        self._current_preview_path = file_path
        self._preview_title.setText(path.name)
        suffix = path.suffix.lower()

        if suffix in {".txt", ".md", ".py", ".json", ".csv", ".log", ".sql", ".yaml", ".yml"}:
            try:
                text = path.read_text(encoding="utf-8", errors="ignore")[:10000]
                self._preview_text.setPlainText(text)
            except Exception:
                self._preview_text.setPlainText("无法读取文件内容。")

        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n".join(paragraphs[:200])
                self._preview_text.setPlainText(text or "(文档内容为空)")
            except ImportError:
                self._preview_text.setPlainText("需要安装 python-docx 库才能预览 .docx 文件。\npip install python-docx")
            except Exception as e:
                self._preview_text.setPlainText(f"无法解析 .docx 文件：{e}")

        elif suffix == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(path), read_only=True, data_only=True)
                lines: list[str] = []
                for sheet_name in wb.sheetnames[:3]:
                    ws = wb[sheet_name]
                    lines.append(f"=== 工作表：{sheet_name} ===")
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 20:
                            lines.append("  … (更多行省略)")
                            break
                        row_str = "\t".join(str(c) if c is not None else "" for c in row)
                        lines.append(row_str)
                    lines.append("")
                wb.close()
                self._preview_text.setPlainText("\n".join(lines))
            except ImportError:
                self._preview_text.setPlainText("需要安装 openpyxl 库才能预览 .xlsx 文件。\npip install openpyxl")
            except Exception as e:
                self._preview_text.setPlainText(f"无法解析 .xlsx 文件：{e}")

        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                lines: list[str] = []
                for i, slide in enumerate(prs.slides):
                    title = ""
                    if slide.shapes.title:
                        title = slide.shapes.title.text
                    lines.append(f"幻灯片 {i + 1}: {title or '(无标题)'}")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                txt = para.text.strip()
                                if txt and txt != title:
                                    lines.append(f"  {txt}")
                self._preview_text.setPlainText("\n".join(lines) or "(演示文稿内容为空)")
            except ImportError:
                self._preview_text.setPlainText("需要安装 python-pptx 库才能预览 .pptx 文件。\npip install python-pptx")
            except Exception as e:
                self._preview_text.setPlainText(f"无法解析 .pptx 文件：{e}")

        else:
            self._preview_text.setPlainText(
                f"二进制文件：{suffix}\n大小：{self._format_size(path)}\n\n点击「用默认程序打开」查看。"
            )

        self._switch_tab("preview")

    # -- files -----------------------------------------------------------------

    def set_workspace(self, path: str, *, refresh: bool = True) -> None:
        from ui.i18n import t

        self._workspace = path or ""
        if path and Path(path).exists():
            self._workspace_label.setText(path)
        else:
            self._workspace_label.setText(t("workspace_unset"))
        self._files_stale = True
        if refresh:
            self.refresh_files()

    def refresh_files(self) -> None:
        self._files_stale = False
        self._file_tree.clear()
        if not self._workspace or not Path(self._workspace).exists():
            from core.settings_runtime import get_workspace_path
            ws = get_workspace_path()
            if ws:
                self._workspace = ws
                self._workspace_label.setText(ws)
            else:
                exports = exports_dir()
                if exports.exists():
                    self._workspace = str(exports)
                    self._workspace_label.setText(str(exports))
                else:
                    return
        root = Path(self._workspace)
        try:
            for item in sorted(root.iterdir()):
                if item.name in _FILE_IGNORED:
                    continue
                tree_item = self._make_file_tree_item(item)
                self._file_tree.addTopLevelItem(tree_item)
        except PermissionError:
            pass

    def _make_file_tree_item(self, path: Path) -> QTreeWidgetItem:
        tree_item = QTreeWidgetItem([path.name, self._format_size(path) if path.is_file() else ""])
        tree_item.setData(0, _FILE_PATH_ROLE, str(path))
        if path.is_dir():
            tree_item.setChildIndicatorPolicy(QTreeWidgetItem.ChildIndicatorPolicy.ShowIndicator)
            tree_item.setData(0, _FILE_LOADED_ROLE, False)
        return tree_item

    def _on_file_tree_expanded(self, item: QTreeWidgetItem) -> None:
        if item.data(0, _FILE_LOADED_ROLE):
            return
        path_str = item.data(0, _FILE_PATH_ROLE)
        if not path_str:
            return
        path = Path(path_str)
        if not path.is_dir():
            return
        item.setData(0, _FILE_LOADED_ROLE, True)
        self._populate_tree_children(item, path)

    def _populate_tree_children(self, parent: QTreeWidgetItem, path: Path) -> None:
        try:
            for item in sorted(path.iterdir()):
                if item.name in _FILE_IGNORED:
                    continue
                parent.addChild(self._make_file_tree_item(item))
        except PermissionError:
            pass

    @staticmethod
    def _format_size(path: Path) -> str:
        try:
            size = path.stat().st_size
        except OSError:
            return ""
        if size < 1024:
            return f"{size} B"
        if size < 1024 * 1024:
            return f"{size / 1024:.1f} KB"
        return f"{size / 1024 / 1024:.1f} MB"

    def _open_file(self, item: QTreeWidgetItem, column: int) -> None:
        path = item.data(0, _FILE_PATH_ROLE)
        if path and Path(path).is_file():
            self._open_file_path(path, self)

    def _iter_file_tree_items(self, parent: QTreeWidgetItem | None = None):
        if parent is None:
            for i in range(self._file_tree.topLevelItemCount()):
                item = self._file_tree.topLevelItem(i)
                yield from self._iter_file_tree_items(item)
            return
        yield parent
        for i in range(parent.childCount()):
            yield from self._iter_file_tree_items(parent.child(i))

    def _selected_file_paths(self) -> list[str]:
        paths: list[str] = []
        for item in self._file_tree.selectedItems():
            path_str = item.data(0, _FILE_PATH_ROLE)
            if path_str and Path(path_str).is_file():
                paths.append(path_str)
        return paths

    def _update_file_batch_bar(self) -> None:
        self._file_batch_bar.set_count(len(self._selected_file_paths()))

    def _file_select_all(self) -> None:
        self._file_tree.clearSelection()
        for item in self._iter_file_tree_items():
            path_str = item.data(0, _FILE_PATH_ROLE)
            if path_str and Path(path_str).is_file():
                item.setSelected(True)
        self._update_file_batch_bar()

    def _file_clear_selection(self) -> None:
        self._file_tree.clearSelection()
        self._update_file_batch_bar()

    def _file_batch_open(self) -> None:
        for path in self._selected_file_paths():
            open_file_path(path, self)

    def _file_batch_delete(self) -> None:
        paths = self._selected_file_paths()
        if not paths:
            return
        reply = QMessageBox.question(
            self,
            "确认删除",
            f"确定删除选中的 {len(paths)} 个文件吗？\n此操作不可恢复。",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return
        for path in paths:
            p = Path(path)
            if p.is_file():
                try:
                    p.unlink()
                except OSError:
                    pass
                remove_artifact_by_path(path)
        self.refresh_file_views()

    def _on_file_tree_context_menu(self, pos) -> None:
        item = self._file_tree.itemAt(pos)
        if not item:
            return
        path_str = item.data(0, _FILE_PATH_ROLE)
        if not path_str:
            return
        path = Path(path_str)
        exec_file_context_menu(
            self,
            path_str,
            self._file_tree.mapToGlobal(pos),
            on_preview=(lambda: self._preview_file(item, 0)) if path.is_file() else None,
            on_after_delete=self.refresh_file_views,
        )

    def _on_file_tree_clicked(self, item: QTreeWidgetItem, column: int) -> None:
        path_str = item.data(0, _FILE_PATH_ROLE)
        if not path_str:
            return
        path = Path(path_str)
        if path.is_file():
            self._file_tab_bar.add_tab(path_str, path.name)
            self._open_file_path(path_str, self)
            return
        self._preview_file(item, column)

    def _on_file_tab_clicked(self, file_path: str) -> None:
        self._current_preview_path = file_path
        path = Path(file_path)
        self._preview_title.setText(path.name)
        suffix = path.suffix.lower()
        if path.is_file() and suffix in {".txt", ".md", ".py", ".json", ".csv", ".log", ".sql", ".yaml", ".yml"}:
            try:
                text = path.read_text(encoding="utf-8", errors="ignore")[:10000]
                self._preview_text.setPlainText(text)
            except Exception:
                self._preview_text.setPlainText("无法读取文件内容。")
        elif path.is_file():
            self._preview_text.setPlainText(
                f"二进制文件：{suffix}\n大小：{self._format_size(path)}\n\n双击文件列表可用默认程序打开。"
            )
        self._switch_tab("preview")

    def _preview_file(self, item: QTreeWidgetItem, column: int) -> None:
        path_str = item.data(0, _FILE_PATH_ROLE)
        if not path_str:
            return
        path = Path(path_str)
        self._current_preview_path = path_str
        self._preview_title.setText(path.name)
        if path.is_file() and path.suffix.lower() in {".txt", ".md", ".py", ".json", ".csv", ".log", ".sql", ".yaml", ".yml"}:
            try:
                text = path.read_text(encoding="utf-8", errors="ignore")[:10000]
                self._preview_text.setPlainText(text)
            except Exception:
                self._preview_text.setPlainText("无法读取文件内容。")
        elif path.is_file():
            self._preview_text.setPlainText(f"二进制文件：{path.suffix}\n大小：{self._format_size(path)}\n\n双击文件列表可用默认程序打开。")
        self._switch_tab("preview")

    def _open_previewed_file(self) -> None:
        if hasattr(self, "_current_preview_path") and self._current_preview_path:
            self._open_file_path(self._current_preview_path, self)

    # -- changes ---------------------------------------------------------------

    def _refresh_changes(self) -> None:
        while self._changes_layout.count() > 1:
            item = self._changes_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()
        task_id = self._task_id
        if not task_id and self._conversation_id:
            from core.task_tracker import latest_task_for_conversation
            task = latest_task_for_conversation(self._conversation_id)
            task_id = task["id"] if task else None
        if not task_id:
            empty = QLabel("暂无变更记录。执行带工具的任务后会在此显示步骤。")
            empty.setObjectName("MutedLabel")
            empty.setWordWrap(True)
            self._changes_layout.insertWidget(0, empty)
            return
        steps = query_all(
            "SELECT * FROM task_steps WHERE task_id=? ORDER BY step_index",
            (task_id,),
        )
        for step in steps:
            card = self._build_change_card(step)
            self._changes_layout.insertWidget(self._changes_layout.count() - 1, card)

    def _build_change_card(self, step: dict) -> QFrame:
        card = QFrame()
        card.setObjectName("ChangeCard")
        card.setStyleSheet(
            "#ChangeCard { background: rgba(255,255,255,0.04); border: 1px solid rgba(255,255,255,0.08); "
            "border-radius: 8px; margin-bottom: 4px; }"
        )
        layout = QVBoxLayout(card)
        layout.setContentsMargins(12, 10, 12, 10)
        layout.setSpacing(6)

        header = QHBoxLayout()
        step_name = step.get("step_name", "未知步骤")
        tool_name = step.get("tool_name", "")
        title = QLabel(f"<b>● {step_name}</b>")
        title.setObjectName("ChangeCardTitle")
        title.setWordWrap(True)
        header.addWidget(title, 1)
        if tool_name:
            tool_label = QLabel(f"🔧 {tool_name}")
            tool_label.setObjectName("MutedLabel")
            header.addWidget(tool_label)
        layout.addLayout(header)

        input_params = step.get("input_json", "") or step.get("input_params", "") or step.get("input_data", "")
        if input_params:
            if isinstance(input_params, dict):
                input_params = json.dumps(input_params, ensure_ascii=False, indent=2)
            input_str = str(input_params)[:300]
            if len(str(input_params)) > 300:
                input_str += " …"
            inp_label = QLabel(f"<span style='color:#888;'>输入：</span>{input_str}")
            inp_label.setObjectName("ChangeCardDetail")
            inp_label.setWordWrap(True)
            layout.addWidget(inp_label)

        output_data = step.get("output_json", "") or step.get("output_data", "") or step.get("result", "")
        if output_data:
            if isinstance(output_data, dict):
                output_data = json.dumps(output_data, ensure_ascii=False, indent=2)
            output_str = str(output_data)[:300]
            if len(str(output_data)) > 300:
                output_str += " …"
            out_label = QLabel(f"<span style='color:#888;'>输出：</span>{output_str}")
            out_label.setObjectName("ChangeCardDetail")
            out_label.setWordWrap(True)
            layout.addWidget(out_label)

        return card

    # -- public helpers --------------------------------------------------------

    def add_artifact(self, art_info: dict) -> None:
        self.refresh_file_views()

    def clear(self) -> None:
        self._task_id = None
        while self._artifact_layout.count() > 1:
            item = self._artifact_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()
        while self._changes_layout.count() > 1:
            item = self._changes_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()
        self._preview_text.clear()
        self._preview_title.setText("选择文件以预览")
